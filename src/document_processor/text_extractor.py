"""
Text extractor module.

This module extracts text from various document formats including PDF, DOCX, PPTX.
It serves as a dispatcher to specific extraction methods based on document type.
"""

import logging
import os
from pathlib import Path
from typing import Dict, List, Union

import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation

from . import ocr
from ..config import SUPPORTED_FILE_TYPES

logger = logging.getLogger(__name__)

class TextExtractor:
    """
    Extracts text from various document formats.
    """
    
    def extract_text(self, file_info: Dict) -> List[Dict]:
        """
        Extract text from a document based on its type.
        
        Args:
            file_info: Dictionary containing file information
                {
                    "path": str,
                    "name": str,
                    "type": str,
                    "metadata": Dict
                }
                
        Returns:
            List[Dict]: List of extracted text blocks with metadata
        """
        file_path = file_info.get("path")
        file_type = file_info.get("type")
        metadata = file_info.get("metadata", {})
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
            
        logger.info(f"Extracting text from {file_path} of type {file_type}")
        
        # Dispatch to appropriate extractor based on file type
        if file_type == "pdf":
            text_blocks = self._extract_from_pdf(file_path, metadata)
        elif file_type == "word":
            text_blocks = self._extract_from_docx(file_path, metadata)
        elif file_type == "powerpoint":
            text_blocks = self._extract_from_pptx(file_path, metadata)
        elif file_type == "image":
            text_blocks = self._extract_from_image(file_path, metadata)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
            
        logger.info(f"Extracted {len(text_blocks)} text blocks from {file_path}")
        return text_blocks
    
    def _extract_from_pdf(self, file_path: str, metadata: Dict) -> List[Dict]:
        """
        Extract text from a PDF file.
        
        Args:
            file_path: Path to the PDF file
            metadata: Document metadata
            
        Returns:
            List[Dict]: List of extracted text blocks with metadata
        """
        text_blocks = []
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                total_pages = len(pdf_reader.pages)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    
                    # Skip empty pages
                    if not text.strip():
                        continue
                        
                    # Add page text with metadata
                    text_blocks.append({
                        "text": text,
                        "metadata": {
                            **metadata,
                            "page": page_num + 1,
                            "total_pages": total_pages
                        }
                    })
        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {e}")
            raise
            
        return text_blocks
    
    def _extract_from_docx(self, file_path: str, metadata: Dict) -> List[Dict]:
        """
        Extract text from a DOCX file.
        
        Args:
            file_path: Path to the DOCX file
            metadata: Document metadata
            
        Returns:
            List[Dict]: List of extracted text blocks with metadata
        """
        text_blocks = []
        
        try:
            doc = DocxDocument(file_path)
            
            # Process paragraphs
            for i, para in enumerate(doc.paragraphs):
                text = para.text.strip()
                if text:
                    text_blocks.append({
                        "text": text,
                        "metadata": {
                            **metadata,
                            "paragraph": i + 1,
                            "section": "body"
                        }
                    })
                    
            # Process tables
            for table_num, table in enumerate(doc.tables):
                table_text = ""
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        table_text += row_text + "\n"
                        
                if table_text.strip():
                    text_blocks.append({
                        "text": table_text.strip(),
                        "metadata": {
                            **metadata,
                            "table": table_num + 1,
                            "section": "table"
                        }
                    })
        except Exception as e:
            logger.error(f"Error extracting text from DOCX {file_path}: {e}")
            raise
            
        return text_blocks
    
    def _extract_from_pptx(self, file_path: str, metadata: Dict) -> List[Dict]:
        """
        Extract text from a PPTX file.
        
        Args:
            file_path: Path to the PPTX file
            metadata: Document metadata
            
        Returns:
            List[Dict]: List of extracted text blocks with metadata
        """
        text_blocks = []
        
        try:
            presentation = Presentation(file_path)
            total_slides = len(presentation.slides)
            
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = ""
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text.strip() + "\n"
                        
                if slide_text.strip():
                    text_blocks.append({
                        "text": slide_text.strip(),
                        "metadata": {
                            **metadata,
                            "slide": slide_num + 1,
                            "total_slides": total_slides
                        }
                    })
        except Exception as e:
            logger.error(f"Error extracting text from PPTX {file_path}: {e}")
            raise
            
        return text_blocks
    
    def _extract_from_image(self, file_path: str, metadata: Dict) -> List[Dict]:
        """
        Extract text from an image using OCR.
        
        Args:
            file_path: Path to the image file
            metadata: Document metadata
            
        Returns:
            List[Dict]: List of extracted text blocks with metadata
        """
        # Delegate to OCR module
        return ocr.extract_text_from_image(file_path, metadata)